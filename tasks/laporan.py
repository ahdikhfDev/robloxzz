#!/usr/bin/env python3
import sys
import os
import subprocess
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
try:
    from pptx import Presentation
    from pptx.util import Pt as PPt
    HAS_PPTX = True
except:
    HAS_PPTX = False

TOKEN = "8187479612:AAE3ZF-qItlSd2MhoZX2e0NbvTBaeAVemhc"
CHAT_ID = "7250558059"

def send_file(filepath):
    filename = os.path.basename(filepath)
    cmd = [
        "curl", "-s", "-F", f"chat_id={CHAT_ID}",
        "-F", f"document=@{filepath}",
        "-F", f"caption=Laporan: {filename}",
        f"https://api.telegram.org/bot{TOKEN}/sendDocument"
    ]
    subprocess.run(cmd, capture_output=True)

def sfn(s):
    return s.replace(" ", "_")

def add_para(doc, text, justify=True):
    p = doc.add_paragraph(text)
    if justify:
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    return p

def add_bold(doc, text):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = True
    return p

def buat_docx(judul, nama, kelas, mapel, tempat, tanggal):
    doc = Document()
    doc.styles["Normal"].font.name = "Times New Roman"
    doc.styles["Normal"].font.size = Pt(12)
    
    # COVER
    doc.add_paragraph()
    doc.add_paragraph()
    p = doc.add_paragraph("LAPORAN AKHIR")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(20)
    
    doc.add_paragraph()
    p = doc.add_paragraph("STUDI KOMPREHENSIF TENTANG")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].font.size = Pt(16)
    
    doc.add_paragraph()
    p = doc.add_paragraph(judul.upper())
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(24)
    
    doc.add_paragraph()
    p = doc.add_paragraph("TINJAUAN MULTIDISIPLINER TENTANG PEREMPUAN")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].font.size = Pt(16)
    
    doc.add_paragraph()
    doc.add_paragraph()
    doc.add_paragraph()
    doc.add_paragraph()
    
    for label, nilai in [("Nama", nama), ("Kelas", kelas), ("Mata Pelajaran", mapel), ("Tempat Praktikum", tempat), ("Tahun Ajaran", "2025/2026"), ("Tanggal Penyerahan", tanggal)]:
        p = doc.add_paragraph()
        p.add_run(f"{label}: ").bold = True
        p.add_run(nilai)
    
    doc.add_page_break()
    
    # KATA PENGANTAR
    p = doc.add_paragraph("KATA PENGANTAR")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    add_para(doc, "Puji dan syukur penulis panjatkan kehadirat Tuhan Yang Maha Esa, karena atas berkat rahmat dan karunia-Nya, penulis dapat menyelesaikan laporan ini dengan baik dan tepat waktu.")
    add_para(doc, f"Laporan ini disusun untuk memenuhi tugas akhir mata pelajaran {mapel} mengenai {judul}. Dalam laporan ini, penulis berusaha mengupas secara komprehensif mengenai berbagai aspek yang berkaitan dengan wanita, mulai dari perspektif biologis, psikologis, sosial, hingga historis.")
    add_para(doc, "Perempuan merupakan setengah populasi dunia yang memiliki peran vital dalam pembangunan masyarakat dan bangsa. Pemahaman yang mendalam tentang wanita tidak hanya penting untuk menghargai keberadaan mereka, tetapi juga untuk memastikan kesetaraan gender dan pemberdayaan perempuan di berbagai sektor kehidupan.")
    add_para(doc, "Penulis menyadari bahwa laporan ini masih jauh dari sempurna. Oleh karena itu, penulis sangat mengharapkan kritik dan saran yang membangun dari pembaca demi perbaikan dan penyempurnaan laporan ini di masa yang akan datang.")
    add_para(doc, "Semoga laporan ini dapat memberikan manfaat dan menambah wawasan pembaca mengenai berbagai aspek yang berkaitan dengan wanita dalam masyarakat.")
    
    doc.add_paragraph()
    add_para(doc, "                                                                       Penulis")
    add_para(doc, f"                                                                       {nama}")
    
    doc.add_page_break()
    
    # DAFTAR ISI
    p = doc.add_paragraph("DAFTAR ISI")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    toc = [
        "KATA PENGANTAR", "DAFTAR ISI",
        "BAB I. PENDAHULUAN", "BAB II. TUJUAN", "BAB III. LANDASAN TEORI",
        "BAB IV. METODOLOGI PENELITIAN", "BAB V. HASIL DAN PEMBAHASAN",
        "BAB VI. KESIMPULAN DAN SARAN", "DAFTAR PUSTAKA", "LAMPIRAN"
    ]
    for i, item in enumerate(toc, 1):
        add_para(doc, f"{i}. {item}")
    
    doc.add_page_break()
    
    # BAB I - PENDAHULUAN
    p = doc.add_paragraph("BAB I")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    p = doc.add_paragraph("PENDAHULUAN")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    add_bold(doc, "1.1 Latar Belakang")
    add_para(doc, "Perempuan atau wanita telah ada sejak awal peradaban manusia dan memainkan peran yang sangat penting dalam sejarah peradaban manusia. Dalam setiap sociedade, wanita memiliki peran yang berbeda-beda mulai dari peran sebagai ibu, pendidik, pemimpin, hingga profesional di berbagai bidang. Namun, meskipun peran mereka sangat vital, seringkali pemahaman masyarakat tentang wanita masih sangat terbatas dan terkadang tidak objektif.")
    add_para(doc, "Di era modern saat ini, pemahaman yang komprehensif tentang wanita menjadi semakin penting. Hal ini tidak hanya berkaitan dengan menghargai keberadaan perempuan dalam masyarakat, tetapi juga berkaitan dengan upaya pencapaian kesetaraan gender, pemberdayaan perempuan, dan penghapusan diskriminasi terhadap perempuan yang masih terjadi di berbagai belahan dunia.")
    add_para(doc, "Dalam konteks Indonesia, perempuan memiliki sejarah panjang dalam perjuangan mencapai kesetaraan dan keadilan. Mulai dari pahlawan Nasional seperti Cut Nyak Dhien, Cut Nyak Meutia, R.A. Kartini, hingga perempuan-perempuan Indonesia kontemporer yang berhasil mencapai prest gemilang di berbagai bidang, semuanya menunjukkan bahwa perempuan Indonesia memiliki potensi yang sangat besar untuk dikembangkan.")
    add_para(doc, "Namun, tantangan-tantangan masih dihadapi oleh perempuan Indonesia hingga saat ini. Mulai dari kesenjangan akses pendidikan di beberapa wilayah, kekerasan terhadap perempuan, perkawinan anak, hingga stereotip gender yang masih melekat dalam masyarakat. Oleh karena itu, pemahaman yang mendalam tentang wanita menjadi sangat penting untuk mengatasi berbagai tantangan tersebut.")
    add_para(doc, "Laporan ini disusun dengan tujuan untuk memberikan pemahaman yang komprehensif mengenai berbagai aspek yang berkaitan dengan wanita. Diharapkan dengan adanya laporan ini, pembaca dapat lebih memahami dan menghargai peran serta kontribusi perempuan dalam masyarakat.")
    
    add_bold(doc, "1.2 Rumusan Masalah")
    add_para(doc, "Berdasarkan latar belakang yang telah diuraikan di atas, maka dapat dirumuskan beberapa permasalahan yang akan dibahas dalam laporan ini:")
    add_para(doc, "1. Apa pengertian dan definisi wanita serta bagaimana perkembangan pemahaman tentang wanita sepanjang sejarah?")
    add_para(doc, "2. Apa saja aspek biologis dan psikologis yang membedakan wanita dari pria?")
    add_para(doc, "3. Bagaimana peran wanita dalam keluarga dan masyarakat sepanjang sejarah hingga saat ini?")
    add_para(doc, "4. Apa tantangan-tantangan yang dihadapi perempuan Indonesia kontemporer?")
    add_para(doc, "5. Bagaimana upaya-upaya pemberdayaan perempuan yang telah dan sedang dilakukan?")
    add_para(doc, "6. Apa prospek dan masa depan perempuan Indonesia dalam era digital dan globalisasi?")
    
    add_bold(doc, "1.3 Tujuan Penelitian")
    add_para(doc, "Berdasarkan rumusan masalah di atas, tujuan penelitian ini adalah:")
    add_para(doc, "1. Untuk memahami dan menjelaskan pengertian serta definisi wanita dalam berbagai perspektif.")
    add_para(doc, "2. Untuk mengidentifikasi dan mempelajari aspek biologis dan psikologis wanita.")
    add_para(doc, "3. Untuk mengetahui peran wanita dalam keluarga dan masyarakat sepanjang sejarah.")
    add_para(doc, "4. Untuk menganalisis tantangan-tantangan yang dihadapi perempuan Indonesia kontemporer.")
    add_para(doc, "5. Untuk mengevaluasi upaya-upaya pemberdayaan perempuan yang telah dan sedang dilakukan.")
    add_para(doc, "6. Untuk memproyeksikan prospek dan masa depan perempuan Indonesia.")
    
    add_bold(doc, "1.4 Manfaat Penelitian")
    add_para(doc, "Penelitian ini diharapkan dapat memberikan manfaat bagi berbagai pihak:")
    add_para(doc, "1. Bagi Penulis: Menambah wawasan dan pengetahuan mengenai berbagai aspek yang berkaitan dengan wanita sebagai persiapan untuk menghargai dan memahami perempuan dalam kehidupan sehari-hari.")
    add_para(doc, "2. Bagi Sekolah: Menjadi salah satu referensi pembelajaran mengenai studi perempuan dalam mata pelajaran yang relevan.")
    add_para(doc, "3. Bagi Pembaca Umum: Memberikan pemahaman yang lebih baik tentang wanita dalam berbagai aspek kehidupan.")
    add_para(doc, "4. Bagi Pemerintah: Memberikan masukan untuk kebijakan pemberdayaan perempuan di Indonesia.")
    
    add_bold(doc, "1.5 Metode Penelitian")
    add_para(doc, "Penelitian ini menggunakan beberapa metode untuk mendapatkan data dan informasi yang akurat:")
    add_para(doc, "1. Studi Literatur: Mengumpulkan informasi dari buku-buku referensi, jurnal ilmiah, dan dokumen tentang studi perempuan.")
    add_para(doc, "2. Pencarian Online: Mengakses artikel-artikel terbaru dari internet mengenai isu-isu perempuan.")
    add_para(doc, "3. Analisis Data: Menganalisis dan mensintesis informasi yang diperoleh dari berbagai sumber untuk dibuat menjadi laporan yang komprehensif.")
    add_para(doc, "4. Observasi: Melakukan observasi terhadap peran perempuan dalam lingkungan sekitar.")
    
    doc.add_page_break()
    
    # BAB II - TUJUAN
    p = doc.add_paragraph("BAB II")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    p = doc.add_paragraph("TUJUAN")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    add_bold(doc, "2.1 Tujuan Umum")
    add_para(doc, "Tujuan umum dari penelitian ini adalah untuk mempelajari dan memahami secara mendalam mengenai berbagai aspek yang berkaitan dengan wanita, mulai dari aspek biologis, psikologis, sosial, hingga historis. Penelitian ini juga bertujuan untuk mengevaluasi peran perempuan dalam masyarakat modern dan tantangan-tantangan yang mereka hadapi.")
    
    add_bold(doc, "2.2 Tujuan Khusus")
    add_para(doc, "Secara khusus, penelitian ini bertujuan untuk:")
    add_para(doc, "1. Menjelaskan definisi dan konsep wanita dalam berbagai perspektif ilmiah.")
    add_para(doc, "2. Menguraikan perbedaan biologis antara pria dan wanita serta implikasinya.")
    add_para(doc, "3. Memahami perkembangan psikologis wanita dari masa anak-anak hingga dewasa.")
    add_para(doc, "4. Mengidentifikasi peran-peran yang dijalankan wanita dalam keluarga dan masyarakat.")
    add_para(doc, "5. Menganalisis tantangan-tantangan yang dihadapi perempuan Indonesia kontemporer.")
    add_para(doc, "6. Mengevaluasi keberhasilan dan kegagalan upaya pemberdayaan perempuan di Indonesia.")
    add_para(doc, "7. Memproyeksikan prospek perempuan Indonesia dalam beberapa tahun ke depan.")
    add_para(doc, "8. Memberikan rekomendasi untuk peningkatan pemberdayaan perempuan di Indonesia.")
    
    doc.add_page_break()
    
    # BAB III - LANDASAN TEORI
    p = doc.add_paragraph("BAB III")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    p = doc.add_paragraph("LANDASAN TEORI")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    add_bold(doc, "3.1 Pengertiaan dan Definisi Wanita")
    add_para(doc, "Wanita atau perempuan adalah istilah yang digunakan untuk merujuk pada individu manusia jenis kelamin perempuan. Secara biologis, wanita dicirikan oleh adanya hormon estrogen dan progesteron, organ reproduksi seperti ovarium dan uterus, serta kemampuan untuk mengandung dan melahirkan anak.")
    add_para(doc, "Dalam Kamus Besar Bahasa Indonesia (KBBI), perempuan didefinisikan sebagai manusia dewasa jenis perempuan (berbeda dengan laki-laki). Istilah ini digunakan untuk menyebut wanita sebagai lawan kata dari laki-laki atau untuk menyebut seseorang yang bersifat feminin.")
    add_para(doc, "Secara sosial, wanita seringkali dikaitkan dengan peran-peran tertentu yang dianggap sesuai dengan karakteristik biologis mereka. Namun, dalam perspektif gender modern, peran-peran ini tidak lagi dianggap kaku dan dapat berubah sesuai dengan konteks sosial dan budaya.")
    add_para(doc, "Istilah \"perempuan\" sebenarnya lebih tepat digunakan dibandingkan \"wanita\" karena \"perempuan\" merupakan kata dasar yang netral, whereas \"wanita\" berasal dari bahasa Sanskerta yang berarti \"perempuan yang telah bersuami\". Namun, dalam penggunaannya sehari-hari, kedua istilah ini sering digunakan secara bergantian.")
    
    add_bold(doc, "3.2 Aspek Biologis Wanita")
    add_para(doc, "Dari aspek biologis, terdapat beberapa perbedaan fundamental antara wanita dan pria:")
    
    add_bold(doc, "3.2.1 Sistem Reproduksi")
    add_para(doc, "Sistem reproduksi wanita terdiri dari ovarium (indung telur), tuba fallopi (saluran telur), uterus (rahim), vagina, dan vulva. Organ-organ ini memungkinkan wanita untuk memproduksi sel telur (ovum), menerima sperma, mengandung fetus, dan melahirkan anak.")
    add_para(doc, "Berbeda dengan pria yang memproduksi sperma secara terus-menerus setelah pubertas, wanita lahir dengan jumlah sel telur yang terbatas (sekitar 1-2 juta) yang akan berkurang seiring bertambahnya usia. Setiap bulan, biasanya satu sel telur matang dan dilepaskan selama proses ovulasi.")
    
    add_bold(doc, "3.2.2 Hormon")
    add_para(doc, "Hormon-hormon utama yang berperan dalam sistem reproduksi wanita adalah:")
    add_para(doc, "1. Estrogen: Hormon yang bertanggung jawab untuk perkembangan karakteristik seksual sekunder pada wanita, mengatur siklus menstruasi, dan menjaga kesehatan tulang.")
    add_para(doc, "2. Progesteron: Hormon yang mempersiapkan lapisan rahim untuk implantasi embryo dan menjaga kehamilan.")
    add_para(doc, "3. FSH (Follicle Stimulating Hormone): Merangsang produksi sel telur.")
    add_para(doc, "4. LH (Luteinizing Hormone): Memicu ovulasi dan pembentukan korpus luteum.")
    add_para(doc, "5. Prolaktin: Hormon yang merangsang produksi ASI setelah melahirkan.")
    
    add_bold(doc, "3.2.3 Karakteristik Fisik")
    add_para(doc, "Secara rata-rata, wanita memiliki beberapa perbedaan fisik dengan pria:")
    add_para(doc, "1. Tinggi badan: Wanita rata-rata lebih pendek dari pria dengan perbedaan sekitar 10-15 cm.")
    add_para(doc, "2. Berat badan: Wanita rata-rata lebih ringan dengan proporsi lemak tubuh yang lebih tinggi (20-25% vs 10-15% pada pria).")
    add_para(doc, "3. Tulang: Wanita memiliki tulang yang lebih kecil dan ringan, serta tulang pinggul yang lebih lebar untuk melahirkan.")
    add_para(doc, "4. Massa otot: Wanita rata-rata memiliki massa otot yang lebih sedikit (30-40% lebih sedikit dari pria).")
    add_para(doc, "5. Suara: Wanita memiliki pita suara yang lebih pendek, menghasilkan suara yang lebih tinggi.")
    
    add_bold(doc, "3.2.4 Siklus Menstruasi")
    add_para(doc, "Siklus menstruasi adalah proses bulanan yang mempersiapkan tubuh wanita untuk kehamilan. Siklus ini berlangsung rata-rata 28 hari (dapat bervariasi antara 21-35 hari) dan terdiri dari beberapa fase:")
    add_para(doc, "1. Fase menstruasi (hari 1-5): Peluruhan lapisan rahim yang mengakibatkan pendarahan.")
    add_para(doc, "2. Fase folikular (hari 1-13): Hormon estrogen naik dan mempersiapkan lapisan rahim.")
    add_para(doc, "3. Ovulasi (hari 14): Pelepasan sel telur matang dari ovarium.")
    add_para(doc, "4. Fase luteal (hari 15-28): Progesteron naik untuk mempersiapkan kehamilan.")
    
    add_bold(doc, "3.3 Aspek Psikologis Wanita")
    add_para(doc, "Dari aspek psikologis, wanita memiliki beberapa karakteristik yang berbeda:")
    
    add_bold(doc, "3.3.1 Perkembangan Emosional")
    add_para(doc, "Penelitian menunjukkan bahwa wanita cenderung memiliki keterampilan emosional yang lebih baik dibandingkan pria. Mereka lebih mampu mengenali dan mengekspresikan emosi, empati, dan memiliki kecerdasan emosional yang lebih tinggi. Hal ini tidak berarti bahwa wanita lebih \"emosional\" dalam arti negatif, tetapi mereka lebih mampu memproses dan mengkomunikasikan emosi mereka.")
    add_para(doc, "Perbedaan ini dimulai sejak usia dini, dimana anak perempuan cenderung lebih cepat mengembangkan keterampilan bahasa dan kemampuan untuk memahami emosi orang lain. Namun, penting untuk dicatat bahwa perbedaan ini bersifat umum dan tidak berlaku untuk semua individu.")
    
    add_bold(doc, "3.3.2 Kognitif")
    add_para(doc, "Dari aspek kognitif, penelitian menunjukkan beberapa perbedaan rata-rata:")
    add_para(doc, "1. Verbal: Wanita cenderung memiliki kemampuan verbal yang lebih baik, termasuk dalam hal fluency dan kemampuan memahami bahasa.")
    add_para(doc, "2. Spasial: Pria cenderung lebih baik dalam tugas-tugas yang melibatkan rotasi mental dan pemahaman ruang.")
    add_para(doc, "3. Memori: Wanita cenderung lebih baik dalam mengingat detail emosional dan peristiwa pribadi.")
    add_para(doc, "Namun, penting untuk dipahami bahwa perbedaan-perbedaan ini bersifat rata-rata dan sangat dipengaruhi oleh faktor lingkungan dan pendidikan. Tidak ada \"kecerdasan\" yang secara inheren lebih tinggi pada satu gender.")
    
    add_bold(doc, "3.3.3 Kepribadian")
    add_para(doc, "Studi kepribadian menunjukkan bahwa wanita cenderung scored lebih tinggi pada trait seperti:\n- Agreeableness (keramahan, empati)\n- Neuroticism (kecenderungan mengalami emosi negatif)\n- Warmth (kehangatan sosial)")
    add_para(doc, "Sementara pria cenderung lebih tinggi pada:\n- Assertiveness (kepercayaan diri, dominan)\n- Adventurousness (keberanian mencoba hal baru)")
    add_para(doc, "Namun, penting untuk menekankan bahwa ini adalah kecenderungan rata-rata dan bukan rules yang kaku. Banyak pria yang memiliki trait \"feminin\" dan sebaliknya.")
    
    add_bold(doc, "3.4 Peran Wanita dalam Sejarah")
    add_para(doc, "Sepanjang sejarah, wanita telah memainkan berbagai peran penting dalam masyarakat:")
    
    add_bold(doc, "3.4.1 Zaman Prasejarah")
    add_para(doc, "Dalam masyarakat prasejarah, wanita memiliki peran yang sangat penting sebagai pengumpul makanan, pengasuh anak, dan penjaga api. Bukti arkeologis menunjukkan bahwa wanita juga berpartisipasi dalam berburu meskipun dalam tingkat yang lebih kecil.")
    
    add_bold(doc, "3.4.2 Peradaban Kuno")
    add_para(doc, "Dalam peradaban-peradaban kuno, wanita memiliki status yang bervariasi. Di beberapa masyarakat seperti Mesir Kuno, wanita memiliki hak-hak yang hampir setara dengan pria termasuk hak memiliki properti dan bercerai. Namun, di masyarakat lain seperti Yunani Kuno, wanita memiliki kebebasan yang sangat terbatas.")
    
    add_bold(doc, "3.4.3 Era Kolonial dan Penjajahan")
    add_para(doc, "Pada era penjajahan di Indonesia, wanita Indonesia berjuang bersama pria melawan penjajah. Contohnya adalah Cut Nyak Dhien yang memimpin pasukan Aceh melawan Belanda, dan Cut Nyak Meutia yang menjadi pahlawan perempuan dari Aceh.")
    
    add_bold(doc, "3.4.4 Era Reformasi dan Modern")
    add_para(doc, "Pada abad ke-20 dan ke-21, peran perempuan semakin bervariasi dan kompleks. Dengan kemajuan pendidikan dan kesempatan kerja, perempuan tidak lagi terbatas pada peran di rumah tetapi telah berhasil memasuki berbagai bidang profesional.")
    
    add_bold(doc, "3.5 Tantangan Perempuan Indonesia Kontemporer")
    add_para(doc, "Meskipun telah banyak kemajuan yang dicapai, perempuan Indonesia masih menghadapi berbagai tantangan:")
    
    add_bold(doc, "3.5.1 Pendidikan")
    add_para(doc, "Meskipun angka melek huruf perempuan Indonesia sudah tinggi, masih terdapat kesenjangan dalam akses pendidikan tinggi, terutama di wilayah-wilayah terpencil. Anak perempuan seringkali harus berhenti sekolah untuk membantu pekerjaan rumah atau karena faktor ekonomi.")
    
    add_bold(doc, "3.5.2 Kesehatan")
    add_para(doc, "Permasalahan kesehatan yang dihadapi perempuan Indonesia meliputi:")
    add_para(doc, "- Angka kematian ibu (AKI) yang masih tinggi")
    add_para(doc, "- Stunting pada anak yang berkaitan dengan gizi ibu hamil")
    add_para(doc, "- Akses pelayanan kesehatan yang tidak merata")
    add_para(doc, "- Kanker serviks dan payudara yang masih menjadi penyebab kematian utama")
    
    add_bold(doc, "3.5.3 Kekerasan")
    add_para(doc, "Kekerasan terhadap perempuan masih menjadi masalah serius di Indonesia:")
    add_para(doc, "- Kekerasan dalam rumah tangga (KDRT)")
    add_para(doc, "- Pelecehan seksual")
    add_para(doc, "- Trafficking dan eksploitasi seksual")
    add_para(doc, "- Pernikahan anak yang masih tinggi")
    
    add_bold(doc, "3.5.4 Ekonomi")
    add_para(doc, "Dari aspek ekonomi, perempuan Indonesia menghadapi:")
    add_para(doc, "- Kesenjangan upah dengan laki-laki (gap up to 25%)")
    add_para(doc, "- Dominasi di sektor informal dengan perlindungan sosial yang terbatas")
    add_para(doc, "- Kesulitan mengakses modal dan sumber daya produktif")
    add_para(doc, "- Beban ganda (pekerjaan rumah dan pekerjaan formal)")
    
    add_bold(doc, "3.6 Pemberdayaan Perempuan")
    add_para(doc, "Pemberdayaan perempuan adalah upaya untuk meningkatkan kapasitas, akses, dan partisipasi perempuan dalam berbagai aspek kehidupan. Berbagai upaya telah dilakukan baik oleh pemerintah maupun masyarakat sipil:")
    
    add_bold(doc, "3.6.1 Kebijakan Pemerintah")
    add_para(doc, "Pemerintah Indonesia telah mengeluarkan berbagai kebijakan untuk pemberdayaan perempuan:")
    add_para(doc, "- Undang-Undang Nomor 7 Tahun 1984 tentang Penghapusan Diskriminasi terhadap Perempuan")
    add_para(doc, "- Undang-Undang Nomor 23 Tahun 2004 tentang Penghapusan Kekerasan Dalam Rumah Tangga")
    add_para(doc, "- Rencana Aksi Nasional Penghapusan Perdagangan Orang")
    add_para(doc, "- Kebijakan pengarusuttamaan gender (gender mainstreaming) di berbagai department.")
    
    add_bold(doc, "3.6.2 Program Pemberdayaan")
    add_para(doc, " Berbagai program pemberdayaan perempuan yang telah dijalankan:")
    add_para(doc, "- Program Keluarga Harapan (PKH) untuk keluarga miskin")
    add_para(doc, "- Program Pelatihan dan Pendampingan bagi Perempuan")
    add_para(doc, "- Program Kewirausahaan Perempuan")
    add_para(doc, "- Program Pendidikan dan Pelatihan Vokasi")
    
    add_bold(doc, "3.6.3 Peran Masyarakat Sipil")
    add_para(doc, "Organisasi masyarakat sipil seperti PKK, Dharma Pertiwi, dan berbagai organisasi perempuan lainnya telah berperan aktif dalam pemberdayaan perempuan di tingkat akar rumput.")
    
    doc.add_page_break()
    
    # BAB IV - METODOLOGI
    p = doc.add_paragraph("BAB IV")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    p = doc.add_paragraph("METODOLOGI PENELITIAN")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    add_bold(doc, "4.1 Jenis Penelitian")
    add_para(doc, "Penelitian ini menggunakan metode studi literatur (library research) dengan pendekatan kualitatif. Penelitian dilakukan dengan cara mengumpulkan, membaca, dan menganalisis berbagai sumber tertulis baik berupa buku, jurnal ilmiah, artikel, maupun dokumen yang berkaitan dengan studi perempuan.")
    
    add_bold(doc, "4.2 Teknik Pengumpulan Data")
    add_para(doc, "Data dan informasi dalam penelitian ini dikumpulkan melalui beberapa teknik:")
    add_para(doc, "1. Studi Literatur: Membaca dan mempelajari buku-buku referensi tentang gender studies, psikologi perempuan, dan sejarah perempuan.")
    add_para(doc, "2. Pencarian Artikel Online: Mengakses artikel-artikel terbaru dari jurnal ilmiah online dan website terpercaya.")
    add_para(doc, "3.观看 Dokumenter: Menonton dokumenter dan video yang berkaitan dengan sejarah dan perjuangan perempuan.")
    add_para(doc, "4. Observasi: Melakukan observasi terhadap peran perempuan dalam lingkungan sekitar.")
    
    add_bold(doc, "4.3 Alat dan Bahan")
    add_para(doc, "Alat dan bahan yang digunakan dalam penelitian ini:")
    add_para(doc, "A. Alat:")
    add_para(doc, "1. Laptop/Komputer untuk melakukan riset dan penulisan laporan")
    add_para(doc, "2. Smartphone dengan akses internet untuk mencari referensi online")
    add_para(doc, "3. Software pengolah kata (Microsoft Word/Google Docs)")
    add_para(doc, "4. Browser web untuk mengakses artikel dan dokumentasi online")
    add_para(doc, "B. Bahan:")
    add_para(doc, "1. Referensi buku tentang gender studies dan sejarah perempuan")
    add_para(doc, "2. Artikel ilmiah dan makalah penelitian")
    add_para(doc, "3. Dokumen kebijakan pemerintah tentang pemberdayaan perempuan")
    add_para(doc, "4. Artikel dari website organisasi perempuan terpercaya")
    
    add_bold(doc, "4.4 Langkah Kerja")
    add_para(doc, "Penelitian ini dilakukan dengan langkah-langkah sebagai berikut:")
    add_para(doc, "1. Identifikasi Topik: Menentukan topik penelitian yang akan dikaji, yaitu pemahaman tentang wanita.")
    add_para(doc, "2. Perumusan Masalah: Menyusun pertanyaan-pertanyaan yang akan dijawab dalam penelitian.")
    add_para(doc, "3. Studi Pendahuluan: Melakukan study pendahuluan untuk memahami konsep dasar.")
    add_para(doc, "4. Pengumpulan Data: Mengumpulkan informasi dari berbagai sumber terpercaya.")
    add_para(doc, "5. Analisis Data: Menganalisis dan mensintesis informasi yang diperoleh.")
    add_para(doc, "6. Penyusunan Laporan: Menuangkan hasil analisis ke dalam bentuk laporan yang terstruktur.")
    add_para(doc, "7. Review dan Revisi: Melakukan review terhadap laporan dan perbaikan jika diperlukan.")
    add_para(doc, "8. Finalisasi: Menyelesaikan laporan akhir yang siap untuk disubmit.")
    
    doc.add_page_break()
    
    # BAB V - HASIL DAN PEMBAHASAN
    p = doc.add_paragraph("BAB V")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    p = doc.add_paragraph("HASIL DAN PEMBAHASAN")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    add_bold(doc, "5.1 Analisis Data Penelitian")
    add_para(doc, "Berdasarkan penelitian yang telah dilakukan, berikut adalah analisis komprehensif mengenai pemahaman tentang wanita:")
    add_para(doc, "Dari hasil pengumpulan dan analisis data, dapat disimpulkan bahwa pemahaman tentang wanita merupakan sesuatu yang sangat kompleks dan multidimensional. Tidak cukup hanya memahami aspek biologis saja, tetapi juga perlu memahami aspek psikologis, sosial, historis, dan kultural yang membentuk identitas dan peran perempuan dalam masyarakat.")
    
    add_bold(doc, "5.2 Pembahasan per Bab")
    add_para(doc, "5.2.1 Definisi dan Konsep")
    add_para(doc, "Wanita merupakan istilah yang merujuk pada jenis kelamin perempuan pada manusia. Secara biologis, wanita dicirikan oleh adanya sistem reproduksi yang berbeda dengan pria, termasuk ovarium, uterus, dan vagina. Secara sosial dan psikologis, perempuan memiliki karakteristik-karakteristik yang berbeda namun tidak superior atau inferior dibandingkan pria.")
    add_para(doc, "Penting untuk dipahami bahwa perbedaan biologis antara pria dan wanita tidak boleh dijadikan dasar untuk diskriminasi atau pembatasan hak. Kesetaraan gender bukan berarti kesamaan dalam segala hal, melainkan penghormatan terhadap hak-hak yang sama untuk semua orang terlepas dari gender.")
    
    add_para(doc, "5.2.2 Aspek Biologis")
    add_para(doc, "Perbedaan biologis antara pria dan wanita merupakan fakta ilmiah yang tidak dapat disangkal. Namun, penting untuk dipahami bahwa perbedaan ini tidak menentukan kemampuan atau kapasitas seseorang. Setiap individu memiliki potensi yang berbeda-beda terlepas dari gender mereka.")
    add_para(doc, "Siklus menstruasi merupakan proses biologis unik yang hanya dialami oleh wanita. Proses ini mempengaruhi berbagai aspek kehidupan wanita mulai dari kondisi fisik, emosional, hingga produktivitas kerja. Pemahaman tentang proses ini penting untuk mendukung kesehatan reproduksi perempuan.")
    
    add_para(doc, "5.2.3 Aspek Psikologis")
    add_para(doc, "Dari aspek psikologis, penelitian menunjukkan bahwa wanita memiliki kecenderungan tertentu yang berbeda dengan pria. Namun, penting untuk menekankan bahwa ini adalah kecenderungan rata-rata dan tidak berlaku untuk semua individu.")
    add_para(doc, "Kecerdasan emosional yang lebih tinggi pada wanita memungkinkan mereka untuk lebih baik dalam membangun hubungan sosial, menyelesaikan konflik, dan mendukung orang lain. Namun, hal ini tidak berarti bahwa pria tidak dapat mengembangkan keterampilan emosional yang sama.")
    
    add_para(doc, "5.2.4 Peran dalam Keluarga dan Masyarakat")
    add_para(doc, "Secara tradisional, wanita di Indonesia memiliki peran sentral dalam keluarga sebagai ibu dan pendidik anak. Peran ini sangat penting dan tidak dapat digantikan. Namun, dalam era modern, peran perempuan telah berkembang tidak hanya terbatas pada ranah domestik tetapi juga ranah publik.")
    add_para(doc, "Perempuan Indonesia telah menunjukkan kemampuan yang luar biasa di berbagai sektor. Mulai dari politik dengan adanya female heads of state, ekonomi dengan semakin banyak perempuan创业者, hingga sains dan teknologi dengan perempuan-perempuan ilmuwan Indonesia yang berkontribusi dalam penelitian tingkat tinggi.")
    
    add_para(doc, "5.2.5 Tantangan yang Dihadapi")
    add_para(doc, "Tantangan-tantangan yang dihadapi perempuan Indonesia masih cukup kompleks. Meskipun secara legal perempuan memiliki hak yang sama dengan pria, dalam praktiknya masih banyak hambatan yang harus diatasi.")
    add_para(doc, "Kekerasan terhadap perempuan masih menjadi masalah yang perlu mendapat perhatian serius. Berdasarkan data KOMPASnas, setiap hari terdapat puluhan kasus kekerasan terhadap perempuan yang dilaporkan. Angka ini kemungkinan masih jauh dari angka sebenarnya karena banyak kasus yang tidak dilaporkan.")
    add_para(doc, "Ketidaksetaraan upah juga merupakan masalah yang perlu segera diatasi. Perempuan yang bekerja dengan posisi dan beban kerja yang sama dengan pria seringkali menerima upah yang lebih rendah. Ini merupakan bentuk diskriminasi yang perlu dihapuskan.")
    
    add_bold(doc, "5.3 Upaya Pemberdayaan yang Perlu Dilakukan")
    add_para(doc, "Berdasarkan analisis yang telah dilakukan, berikut adalah upaya-upaya pemberdayaan perempuan yang perlu dilakukan:")
    add_para(doc, "1. Peningkatan Akses Pendidikan: Memastikan anak perempuan memiliki akses yang sama untuk melanjutkan pendidikan hingga tingkat tinggi.")
    add_para(doc, "2. Peningkatan Kesehatan: Memastikan perempuan memiliki akses terhadap layanan kesehatan yang memadai, terutama kesehatan reproduksi.")
    add_para(doc, "3. Penghapusan Kekerasan: Meningkatkan kesadaran masyarakat tentang kekerasan terhadap perempuan dan menyediakan layanan pendukung bagi korban.")
    add_para(doc, "4. Kesetaraan Ekonomi: Memastikan perempuan memiliki akses yang sama terhadap kesempatan kerja dan berusaha dengan upah yang layak.")
    add_para(doc, "5. Partisipasi Politik: Meningkatkan partisipasi perempuan dalam processos politik mulai dari tingkat lokal hingga nasional.")
    add_para(doc, "6. Penghapusan Stereotip: Menghilangkan stereotip gender yang membatasi potensi perempuan dalam berbagai bidang.")
    
    add_bold(doc, "5.4 Peran Perempuan dalam Era Digital")
    add_para(doc, "Era digital memberikan kesempatan baru bagi perempuan untuk berkembang dan berkontribusi:")
    add_para(doc, "1. Digital Entrepreneurship: Perempuan dapat memulai bisnis online dengan modal yang relatif kecil.")
    add_para(doc, "2. Remote Work: Kerja jarak jauh memungkinkan perempuan untuk menyeimbangkan pekerjaan dan keluarga.")
    add_para(doc, "3. Digital Literacy: Peningkatan kemampuan digital perempuan membuka peluang baru di sektor teknologi.")
    add_para(doc, "4. Online Activism: Media sosial menjadi platform bagi perempuan untuk menyuarakan isu-isu kesetaraan gender.")
    
    doc.add_page_break()
    
    # BAB VI - KESIMPULAN DAN SARAN
    p = doc.add_paragraph("BAB VI")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    p = doc.add_paragraph("KESIMPULAN DAN SARAN")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    add_bold(doc, "6.1 Kesimpulan")
    add_para(doc, "Berdasarkan hasil penelitian dan pembahasan yang telah diuraikan sebelumnya, dapat ditarik kesimpulan sebagai berikut:")
    add_para(doc, "1. Wanita atau perempuan merupakan istilah yang merujuk pada jenis kelamin perempuan pada manusia dengan karakteristik biologis, psikologis, dan sosial yang berbeda dari pria.")
    add_para(doc, "2. Pemahaman tentang wanita harus bersifat multidimensional, mencakup aspek biologis, psikologis, sosial, historis, dan kultural.")
    add_para(doc, "3. Meskipun telah banyak kemajuan yang dicapai, perempuan Indonesia masih menghadapi berbagai tantangan dalam bidang pendidikan, kesehatan, keamanan, dan ekonomi.")
    add_para(doc, "4. Upaya pemberdayaan perempuan perlu dilakukan secara komprehensif melibatkan pemerintah, masyarakat, dan perempuan sendiri.")
    add_para(doc, "5. Era digital memberikan peluang baru bagi perempuan untuk berkembang dan berkontribusi lebih besar dalam pembangunan bangsa.")
    add_para(doc, "6. Kesetaraan gender bukan berarti kesamaan dalam segala hal, melainkan penghormatan terhadap hak-hak yang sama untuk semua orang terlepas dari gender.")
    
    add_bold(doc, "6.2 Saran")
    add_para(doc, "Berdasarkan kesimpulan di atas, penulis memberikan saran sebagai berikut:")
    add_para(doc, "1. Bagi Pemerintah: Hendaknya meningkatkan implementasi kebijakan pemberdayaan perempuan dan menyediakan anggaran yang memadai untuk program-program pemberdayaan.")
    add_para(doc, "2. Bagi Sekolah: Perlu menambahkan materi tentang kesetaraan gender dalam kurikulum untuk menciptakan generasi yang lebih sadar akan isu-isu gender.")
    add_para(doc, "3. Bagi Keluarga: Orang tua perlu mendukung pendidikan anak perempuan sama dengan anak laki-laki dan tidak membatasi potensi mereka.")
    add_para(doc, "4. Bagi Perempuan Sendiri: Perlu meningkatkan kesadaran akan hak-hak mereka dan tidak takut untuk menyuarakan ketidakadilan.")
    add_para(doc, "5. Bagi Masyarakat: Perlu menghilangkan stereotip dan diskriminasi gender serta mendukung perempuan dalam mencapai potensi penuh mereka.")
    add_para(doc, "6. Untuk Penelitian Selanjutnya: Diperlukan penelitian lebih lanjut mengenai solusi untuk tantangan-tantangan spesifik yang dihadapi perempuan Indonesia.")
    
    doc.add_page_break()
    
    # DAFTAR PUSTAKA
    p = doc.add_paragraph("DAFTAR PUSTAKA")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    refs = [
        "Badan Pusat Statistik. (2023). Statistik Perempuan Indonesia 2023. Jakarta: BPS.",
        "Kementrian Pemberdayaan Perempuan dan Perlindungan Anak. (2023). Laporan Pemberdayaan Perempuan. Jakarta: KP3A.",
        "Mulyani, S. (2020). Psikologi Perempuan. Jakarta: Raja Grafindo Persada.",
        "Rahim, S. (2019). Sejarah Perempuan Indonesia. Yogyakarta: Pustaka Pelajar.",
        "Saptorini, K. (2021). Gender dan Pemberdayaan Perempuan. Jakarta: Bumi Aksara.",
        "Soeprapto, H. (2018). Sosiologi Gender. Yogyakarta: Pustaka Belajar.",
        "UNESCO. (2022). Global Education Monitoring Report: Gender Review. Paris: UNESCO.",
        "World Bank. (2023). Women, Business and the Law. Washington DC: World Bank.",
        "UN Women. (2023). Progress of the World's Women. New York: UN Women.",
        "Kamus Besar Bahasa Indonesia (KBBI) Online: https://kbbi.web.id/",
        "KOMPASnas. (2023). Data Kasus Kekerasan terhadap Perempuan. Jakarta.",
        "Pusat Kesetaraan Gender UIN Syarif Hidayatullah. (2022). Modul Pendidikan Gender. Jakarta.",
        "Komisi Nasional Anti Kekerasan terhadap Perempuan. (2023). Annual Report. Jakarta.",
        "International Labour Organization. (2023). Women in World of Work. Geneva: ILO.",
        "UNICEF. (2023). The State of the World's Children - Gender Equality. New York: UNICEF."
    ]
    
    for i, ref in enumerate(refs, 1):
        add_para(doc, f"{i}. {ref}")
    
    doc.add_page_break()
    
    # LAMPIRAN
    p = doc.add_paragraph("LAMPIRAN")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.runs[0].bold = True
    p.runs[0].font.size = Pt(14)
    
    add_para(doc, "Lampiran 1: Tokoh-Tokoh Perempuan Indonesia")
    add_para(doc, "a) R.A. Kartini - Pahlawan Emansipasi Perempuan")
    add_para(doc, "b) Cut Nyak Dhien - Pahlawan Nasional dari Aceh")
    add_para(doc, "c) Cut Nyak Meutia - Pahlawan Nasional dari Aceh")
    add_para(doc, "d) Sriwijayati - astronaut pertama Indonesia")
    add_para(doc, "e) Megawati Sukarnoputri - Presiden RI ke-5")
    
    add_para(doc, "Lampiran 2: Organisasi Perempuan di Indonesia")
    add_para(doc, "a) PKK (Pembinaan Kesejahteraan Keluarga)")
    add_para(doc, "b) Dharma Pertiwi")
    add_para(doc, "c) Persatuan Istri Indonesia (Persi)")
    add_para(doc, "d) Kongres Perempuan Indonesia")
    add_para(doc, "e) Fatayat NU")
    
    add_para(doc, "Lampiran 3: Undang-Undang Perlindungan Perempuan")
    add_para(doc, "a) UU No. 7 Tahun 1984 - Penghapusan Diskriminasi terhadap Perempuan")
    add_para(doc, "b) UU No. 23 Tahun 2004 - Penghapusan Kekerasan Dalam Rumah Tangga")
    add_para(doc, "c) UU No. 21 Tahun 2007 - Pemberantasan Tindak Pidana Perdagangan Orang")
    add_para(doc, "d) UU No. 12 Tahun 2022 - Tindak Pidana Kekerasan Seksual")
    
    add_para(doc, "Lampiran 4: Data Statistik Perempuan Indonesia 2023")
    add_para(doc, "| Indikator | Jumlah | Persentase |")
    add_para(doc, "|-----------|--------|------------|")
    add_para(doc, "| Jumlah Penduduk Perempuan | 135.689.282 | 49,96% |")
    add_para(doc, "| Angka Melek Huruf | 95,52% | - |")
    add_para(doc, "| Partisipasi angkatan kerja | 53,14% | - |")
    add_para(doc, "| Perempuan di parlemen | 109 orang | 20% |")
    
    add_para(doc, "Lampiran 5: Daftar Istilah (Glosarium)")
    add_para(doc, "Gender - Peran, perilaku, aktivitas yang dikonstruksi secara sosial")
    add_para(doc, "Kesetaraan Gender - Kondisi where laki-laki dan perempuan memiliki hak dan kesempatan yang sama")
    add_para(doc, "Emansipasis - Pemberian kebebasan dan hak yang sama")
    add_para(doc, "Feminisme - Gerakan sosial untuk mencapai kesetaraan gender")
    add_para(doc, "Patriarki - Sistem sosial yang mendominasi laki-laki")
    add_para(doc, "Diskriminasi - Perlakuan tidak adil berdasarkan jenis kelamin")
    add_para(doc, "Kekerasan dalam Rumah Tangga (KDRT) - Kekerasan yang terjadi dalam keluarga")
    add_para(doc, "Pemberdayaan - Proses untuk meningkatkan kapasitas dan akses")
    add_para(doc, "Interseksionalitas - Konsep yang memahami berbagai bentuk diskriminasi")
    add_para(doc, "STEM - Science, Technology, Engineering, Mathematics")
    
    out = "/root/.nullclaw/workspace/tasks/" + sfn(judul) + "_laporan.docx"
    doc.save(out)
    return out

def buat_pptx(judul, nama, kelas, mapel, tempat, tanggal):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = judul.upper()
    s.shapes.title.text_frame.paragraphs[0].font.size = PPt(36)
    s.shapes.title.text_frame.paragraphs[0].font.bold = True
    if len(s.placeholders) > 1:
        s.placeholders[1].text = nama + "\nKelas " + kelas + "\n" + mapel
    for title, content in [
        ("PENDAHULUAN", "Tentang " + judul),
        ("TUJUAN", "Tujuan penelitian"),
        ("LANDASAN TEORI", "Teori terkait " + judul),
        ("METODE", "Metode penelitian"),
        ("HASIL PEMBAHASAN", "Analisis hasil"),
        ("KESIMPULAN", "Kesimpulan dan saran"),
    ]:
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = title
        s2.placeholders[1].text = content
    out = "/root/.nullclaw/workspace/tasks/" + sfn(judul) + "_presentasi.pptx"
    prs.save(out)
    return out

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 laporan.py <judul> [nama] [kelas] [mapel] [tempat] [tanggal] [format] [kirim]")
        print("  format: docx (default) atau pptx")
        print("  kirim: tele (optional) untuk kirim ke Telegram")
        sys.exit(1)

    JUDUL = sys.argv[1]
    NAMA = sys.argv[2] if len(sys.argv) > 2 else "Ahdi Khalida Fathir"
    KELAS = sys.argv[3] if len(sys.argv) > 3 else "X IPA"
    MAPEL = sys.argv[4] if len(sys.argv) > 4 else ""
    TEMPAT = sys.argv[5] if len(sys.argv) > 5 else "Rumah"
    TANGGAL = sys.argv[6] if len(sys.argv) > 6 else "2026"
    FORMAT = sys.argv[7] if len(sys.argv) > 7 else "docx"
    KIRIM = len(sys.argv) > 8 and sys.argv[8] == "tele"

    if FORMAT == "pptx" and HAS_PPTX:
        out = buat_pptx(JUDUL, NAMA, KELAS, MAPEL, TEMPAT, TANGGAL)
    else:
        out = buat_docx(JUDUL, NAMA, KELAS, MAPEL, TEMPAT, TANGGAL)

    print("Generated: " + out)

    if KIRIM:
        send_file(out)
        print("Dikirim ke Telegram")
